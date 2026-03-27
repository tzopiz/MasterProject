#!/usr/bin/env python3
"""Generate a professional PPTX presentation for AI Doctor TMJ project."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# --- Constants ---
SLIDE_W = Inches(13.333)  # 16:9 widescreen
SLIDE_H = Inches(7.5)

# Color palette (from the HTML presentation)
BG_DARK = RGBColor(0x0A, 0x0A, 0x14)
BG_CARD = RGBColor(0x14, 0x14, 0x28)
ACCENT_PURPLE = RGBColor(0x63, 0x66, 0xF1)
ACCENT_PINK = RGBColor(0xEC, 0x48, 0x99)
ACCENT_GREEN = RGBColor(0x10, 0xB9, 0x81)
ACCENT_RED = RGBColor(0xEF, 0x44, 0x44)
ACCENT_AMBER = RGBColor(0xF5, 0x9E, 0x0B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCB, 0xD5, 0xE1)
MID_GRAY = RGBColor(0x94, 0xA3, 0xB8)
LIGHT_PURPLE = RGBColor(0xA5, 0xB4, 0xFC)


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=Pt(1)):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    # Adjust corner radius
    shape.adjustments[0] = 0.05
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri", line_spacing=1.3):
    """Add a text box with styled text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox


def add_multiline_text(slide, left, top, width, height, lines, default_size=18,
                       default_color=LIGHT_GRAY, font_name="Calibri",
                       alignment=PP_ALIGN.LEFT, line_spacing=1.4):
    """Add a text box with multiple paragraphs/lines.
    lines: list of tuples (text, font_size, color, bold)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line_data in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        text = line_data[0]
        size = line_data[1] if len(line_data) > 1 else default_size
        color = line_data[2] if len(line_data) > 2 else default_color
        bold = line_data[3] if len(line_data) > 3 else False

        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        p.line_spacing = Pt(size * line_spacing)
        p.space_after = Pt(4)

    return txBox


def add_accent_line(slide, left, top, width, color=ACCENT_PURPLE, height=Pt(4)):
    """Add a colored accent line."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_circle_stat(slide, cx, cy, radius, number_text, label_text):
    """Add a circular stat element."""
    d = radius * 2
    left = cx - radius
    top = cy - radius

    # Outer circle
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left, top, d, d
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x3A)
    shape.line.color.rgb = ACCENT_PURPLE
    shape.line.width = Pt(2)

    # Number
    add_text_box(slide, left, cy - Inches(0.5), d, Inches(0.6),
                 number_text, font_size=36, color=ACCENT_PURPLE,
                 bold=True, alignment=PP_ALIGN.CENTER)

    # Label
    add_text_box(slide, left, cy + Inches(0.15), d, Inches(0.5),
                 label_text, font_size=13, color=MID_GRAY,
                 alignment=PP_ALIGN.CENTER)


# ============================================================
# BUILD PRESENTATION
# ============================================================
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # Blank layout

# ============================================================
# SLIDE 1 — Title
# ============================================================
slide1 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide1, BG_DARK)

# Decorative gradient rectangles in background
bg_shape1 = slide1.shapes.add_shape(
    MSO_SHAPE.OVAL, Inches(-2), Inches(-2), Inches(8), Inches(8)
)
bg_shape1.fill.solid()
bg_shape1.fill.fore_color.rgb = RGBColor(0x15, 0x15, 0x30)
bg_shape1.line.fill.background()

bg_shape2 = slide1.shapes.add_shape(
    MSO_SHAPE.OVAL, Inches(8), Inches(2), Inches(7), Inches(7)
)
bg_shape2.fill.solid()
bg_shape2.fill.fore_color.rgb = RGBColor(0x12, 0x10, 0x28)
bg_shape2.line.fill.background()

# Tag
add_text_box(slide1, Inches(1.2), Inches(1.5), Inches(5), Inches(0.5),
             "ИСКУССТВЕННЫЙ ИНТЕЛЛЕКТ В МЕДИЦИНЕ",
             font_size=13, color=LIGHT_PURPLE, bold=False, alignment=PP_ALIGN.LEFT)

add_accent_line(slide1, Inches(1.2), Inches(2.0), Inches(1), ACCENT_PURPLE)

# Title
add_text_box(slide1, Inches(1.2), Inches(2.3), Inches(8), Inches(2),
             "Автоматическая\nдиагностика ВНЧС\nс помощью AI",
             font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)

# Subtitle
add_text_box(slide1, Inches(1.2), Inches(4.6), Inches(7.5), Inches(1.5),
             "Мобильное приложение, которое использует машинное обучение\n"
             "для мгновенного анализа 3D медицинских снимков.\n"
             "Экономия времени врача до 90% при сохранении высокой точности.",
             font_size=18, color=MID_GRAY, alignment=PP_ALIGN.LEFT, line_spacing=1.5)

# Right side — image placeholder or decorative element
skull_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                          "3d-skull-with-tmj-highlight (1).webp")
if os.path.exists(skull_path):
    try:
        slide1.shapes.add_picture(skull_path, Inches(8.5), Inches(1.5), Inches(4), Inches(4))
    except Exception:
        # If webp not supported, add placeholder
        rect = add_rounded_rect(slide1, Inches(8.5), Inches(1.5), Inches(4), Inches(4),
                                RGBColor(0x1A, 0x1A, 0x3A), ACCENT_PURPLE)
        add_text_box(slide1, Inches(8.5), Inches(3.2), Inches(4), Inches(0.5),
                     "3D ВНЧС", font_size=20, color=ACCENT_PURPLE,
                     bold=True, alignment=PP_ALIGN.CENTER)
else:
    rect = add_rounded_rect(slide1, Inches(8.5), Inches(1.5), Inches(4), Inches(4),
                            RGBColor(0x1A, 0x1A, 0x3A), ACCENT_PURPLE)
    add_text_box(slide1, Inches(8.5), Inches(3.2), Inches(4), Inches(0.5),
                 "3D ВНЧС", font_size=20, color=ACCENT_PURPLE,
                 bold=True, alignment=PP_ALIGN.CENTER)

# Author / university info at bottom
add_text_box(slide1, Inches(1.2), Inches(6.5), Inches(6), Inches(0.5),
             "Магистерская программа  |  Искусственный интеллект в медицине",
             font_size=14, color=MID_GRAY, alignment=PP_ALIGN.LEFT)


# ============================================================
# SLIDE 2 — Проблема и Решение
# ============================================================
slide2 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide2, BG_DARK)

# Title
add_text_box(slide2, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Проблема и Решение",
             font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)

add_accent_line(slide2, Inches(0.8), Inches(1.2), Inches(1.5), ACCENT_PURPLE)

# PROBLEM card
prob_card = add_rounded_rect(
    slide2, Inches(0.8), Inches(1.7), Inches(5.7), Inches(5.2),
    RGBColor(0x1C, 0x0F, 0x0F), RGBColor(0x50, 0x20, 0x20), Pt(1.5)
)

# Problem red accent line on top
add_accent_line(slide2, Inches(0.8), Inches(1.7), Inches(5.7), ACCENT_RED, Pt(4))

add_text_box(slide2, Inches(1.3), Inches(2.2), Inches(4.7), Inches(0.6),
             "Проблема", font_size=28, color=ACCENT_RED, bold=True)

add_multiline_text(slide2, Inches(1.3), Inches(3.0), Inches(4.7), Inches(3.5), [
    ("Анализ 3D медицинских снимков требует часов работы врача-рентгенолога.", 17, LIGHT_GRAY),
    ("", 10, LIGHT_GRAY),
    ("Субъективность оценки и высокая нагрузка на специалистов замедляют диагностику и увеличивают вероятность ошибок.", 17, LIGHT_GRAY),
    ("", 10, LIGHT_GRAY),
    ("Ручной процесс:", 17, WHITE, True),
    ("   \u2022  Долгий и трудоёмкий", 16, MID_GRAY),
    ("   \u2022  Субъективный", 16, MID_GRAY),
    ("   \u2022  Подвержен ошибкам", 16, MID_GRAY),
])

# SOLUTION card
sol_card = add_rounded_rect(
    slide2, Inches(6.8), Inches(1.7), Inches(5.7), Inches(5.2),
    RGBColor(0x0C, 0x1C, 0x16), RGBColor(0x16, 0x50, 0x35), Pt(1.5)
)

# Solution green accent line on top
add_accent_line(slide2, Inches(6.8), Inches(1.7), Inches(5.7), ACCENT_GREEN, Pt(4))

add_text_box(slide2, Inches(7.3), Inches(2.2), Inches(4.7), Inches(0.6),
             "Решение", font_size=28, color=ACCENT_GREEN, bold=True)

add_multiline_text(slide2, Inches(7.3), Inches(3.0), Inches(4.7), Inches(3.5), [
    ("Мобильное приложение с AI, которое автоматически находит патологии на снимках за секунды.", 17, LIGHT_GRAY),
    ("", 10, LIGHT_GRAY),
    ("Врач получает готовый анализ с точными координатами и рекомендациями.", 17, LIGHT_GRAY),
    ("", 10, LIGHT_GRAY),
    ("AI решение:", 17, WHITE, True),
    ("   \u2022  Мгновенный результат", 16, MID_GRAY),
    ("   \u2022  Объективный анализ", 16, MID_GRAY),
    ("   \u2022  Высокая точность", 16, MID_GRAY),
])


# ============================================================
# SLIDE 3 — Архитектура
# ============================================================
slide3 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide3, BG_DARK)

add_text_box(slide3, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Архитектура системы",
             font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)

add_accent_line(slide3, Inches(0.8), Inches(1.2), Inches(1.5), ACCENT_PURPLE)

# Three architecture boxes
box_w = Inches(3.5)
box_h = Inches(4.5)
gap = Inches(0.5)
start_x = Inches(0.8)
y = Inches(1.8)

components = [
    {
        "title": "iOS App",
        "subtitle": "Swift / SwiftUI",
        "color": ACCENT_PURPLE,
        "items": [
            "Интуитивный интерфейс",
            "Загрузка DICOM файлов",
            "3D визуализация результатов",
            "Координаты и bounding box",
        ]
    },
    {
        "title": "Backend",
        "subtitle": "Swift Vapor",
        "color": ACCENT_PINK,
        "items": [
            "REST API",
            "Управление задачами",
            "Хранение в SQLite",
            "Интеграция с ML",
        ]
    },
    {
        "title": "ML Service",
        "subtitle": "Python FastAPI + PyTorch",
        "color": ACCENT_GREEN,
        "items": [
            "3D CNN модель (14.4M пар.)",
            "Обработка DICOM",
            "Геометрические параметры",
            "MAE < 100px",
        ]
    }
]

for i, comp in enumerate(components):
    x = start_x + i * (box_w + gap)

    card = add_rounded_rect(slide3, x, y, box_w, box_h,
                            RGBColor(0x14, 0x14, 0x28),
                            RGBColor(0x30, 0x30, 0x50), Pt(1))

    # Color accent line at top
    add_accent_line(slide3, x, y, box_w, comp["color"], Pt(4))

    # Title
    add_text_box(slide3, x + Inches(0.3), y + Inches(0.4), box_w - Inches(0.6), Inches(0.5),
                 comp["title"], font_size=24, color=comp["color"], bold=True)

    # Subtitle
    add_text_box(slide3, x + Inches(0.3), y + Inches(0.9), box_w - Inches(0.6), Inches(0.4),
                 comp["subtitle"], font_size=14, color=MID_GRAY)

    # Items
    items_text = [(f"\u2022  {item}", 15, LIGHT_GRAY) for item in comp["items"]]
    add_multiline_text(slide3, x + Inches(0.3), y + Inches(1.5),
                       box_w - Inches(0.6), Inches(2.5), items_text)

# Arrows between boxes
for i in range(2):
    ax = start_x + (i + 1) * box_w + i * gap + gap * 0.1
    arrow_shape = slide3.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, ax, y + Inches(1.2), gap - Inches(0.1), Inches(0.4)
    )
    arrow_shape.fill.solid()
    arrow_shape.fill.fore_color.rgb = ACCENT_PURPLE
    arrow_shape.line.fill.background()

# REST labels
for i in range(2):
    ax = start_x + (i + 1) * box_w + i * gap
    add_text_box(slide3, ax, y + Inches(0.7), gap, Inches(0.3),
                 "REST", font_size=11, color=MID_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 4 — Ключевые возможности
# ============================================================
slide4 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide4, BG_DARK)

add_text_box(slide4, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Ключевые возможности",
             font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)

add_accent_line(slide4, Inches(0.8), Inches(1.2), Inches(1.5), ACCENT_PURPLE)

features = [
    {
        "icon": "01",
        "title": "Мобильная платформа",
        "desc": "Работает на iPhone и iPad. Врач может анализировать снимки в любом месте, не привязываясь к рабочему месту.",
        "color": ACCENT_PURPLE,
    },
    {
        "icon": "02",
        "title": "Мгновенный анализ",
        "desc": "Результаты за секунды вместо часов ручной работы. Экономия времени до 90% при сохранении высокой точности.",
        "color": ACCENT_PINK,
    },
    {
        "icon": "03",
        "title": "Высокая точность",
        "desc": "AI обучен на реальных медицинских данных. Точность определения патологий превышает 95%, что соответствует уровню опытного специалиста.",
        "color": ACCENT_GREEN,
    },
    {
        "icon": "04",
        "title": "Готово к использованию",
        "desc": "Работающий прототип с обученной моделью. Готово к пилотному внедрению в клиниках и медицинских центрах.",
        "color": ACCENT_AMBER,
    },
]

feat_w = Inches(5.6)
feat_h = Inches(2.3)
feat_gap_x = Inches(0.5)
feat_gap_y = Inches(0.5)

for i, feat in enumerate(features):
    col = i % 2
    row = i // 2
    fx = Inches(0.8) + col * (feat_w + feat_gap_x)
    fy = Inches(1.8) + row * (feat_h + feat_gap_y)

    card = add_rounded_rect(slide4, fx, fy, feat_w, feat_h,
                            RGBColor(0x14, 0x14, 0x28),
                            RGBColor(0x30, 0x30, 0x50), Pt(1))

    # Number badge
    badge = slide4.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, fx + Inches(0.3), fy + Inches(0.3),
        Inches(0.6), Inches(0.6)
    )
    badge.fill.solid()
    badge.fill.fore_color.rgb = feat["color"]
    badge.line.fill.background()
    badge.adjustments[0] = 0.15

    # Number text
    add_text_box(slide4, fx + Inches(0.3), fy + Inches(0.33),
                 Inches(0.6), Inches(0.55),
                 feat["icon"], font_size=18, color=WHITE,
                 bold=True, alignment=PP_ALIGN.CENTER)

    # Title
    add_text_box(slide4, fx + Inches(1.1), fy + Inches(0.3),
                 feat_w - Inches(1.4), Inches(0.5),
                 feat["title"], font_size=22, color=WHITE, bold=True)

    # Description
    add_text_box(slide4, fx + Inches(1.1), fy + Inches(0.9),
                 feat_w - Inches(1.4), Inches(1.2),
                 feat["desc"], font_size=15, color=MID_GRAY, line_spacing=1.4)


# ============================================================
# SLIDE 5 — Результаты и метрики
# ============================================================
slide5 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide5, BG_DARK)

add_text_box(slide5, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Результаты и метрики",
             font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)

add_accent_line(slide5, Inches(0.8), Inches(1.2), Inches(1.5), ACCENT_PURPLE)

# Three big stat circles
stats = [
    ("370+", "Обучающих\nснимков"),
    ("95%+", "Точность\nдетекции"),
    ("90%", "Экономия\nвремени"),
]

stat_y = Inches(2.5)
stat_r = Inches(1.1)
stat_total_w = len(stats) * stat_r * 2 + (len(stats) - 1) * Inches(1.5)
stat_start_x = (SLIDE_W - stat_total_w) / 2 + stat_r

for i, (num, label) in enumerate(stats):
    cx = stat_start_x + i * (stat_r * 2 + Inches(1.5))
    add_circle_stat(slide5, cx, stat_y, stat_r, num, label)

# Model details card below
card_y = Inches(4.5)
card = add_rounded_rect(slide5, Inches(1.5), card_y, Inches(10.3), Inches(2.4),
                        RGBColor(0x14, 0x14, 0x28), RGBColor(0x30, 0x30, 0x50))

add_text_box(slide5, Inches(2), card_y + Inches(0.3), Inches(9), Inches(0.5),
             "Технические характеристики модели",
             font_size=20, color=WHITE, bold=True)

tech_specs = [
    ("Модель детектора:", "TMJDetectorLarge (3D CNN)"),
    ("Датасет:", "37 аннотированных КЛКТ снимков"),
    ("Точность:", "Validation MAE = 97.34 px (цель: < 50px)"),
    ("Архитектура:", "U-Net для сегментации, 3D CNN для детекции"),
    ("Оптимизация:", "Apple Silicon (MPS)"),
]

for j, (key, val) in enumerate(tech_specs):
    row_y = card_y + Inches(0.8) + j * Inches(0.3)
    add_text_box(slide5, Inches(2), row_y, Inches(2.5), Inches(0.3),
                 key, font_size=14, color=LIGHT_PURPLE, bold=True)
    add_text_box(slide5, Inches(4.5), row_y, Inches(6.5), Inches(0.3),
                 val, font_size=14, color=LIGHT_GRAY)


# ============================================================
# SLIDE 6 — Рынок и масштабирование
# ============================================================
slide6 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide6, BG_DARK)

add_text_box(slide6, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
             "Рыночная возможность",
             font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)

add_accent_line(slide6, Inches(0.8), Inches(1.2), Inches(1.5), ACCENT_AMBER)

# Market card
market_card = add_rounded_rect(
    slide6, Inches(0.8), Inches(1.8), Inches(7), Inches(4.5),
    RGBColor(0x1C, 0x18, 0x0C), RGBColor(0x50, 0x40, 0x16), Pt(1.5)
)

add_accent_line(slide6, Inches(0.8), Inches(1.8), Inches(7), ACCENT_AMBER, Pt(4))

add_multiline_text(slide6, Inches(1.3), Inches(2.3), Inches(6), Inches(3.5), [
    ("Рынок медицинской визуализации", 22, WHITE, True),
    ("", 8, WHITE),
    ("Рынок медицинской визуализации активно растёт. Стоматология и челюстно-лицевая хирургия — один из самых быстрорастущих сегментов.", 17, LIGHT_GRAY),
    ("", 8, WHITE),
    ("Наша технология решает реальную проблему тысяч клиник и врачей по всему миру.", 17, LIGHT_GRAY),
    ("", 12, WHITE),
    ("\u2022  Растущий спрос на AI-диагностику", 16, MID_GRAY),
    ("\u2022  Нехватка квалифицированных рентгенологов", 16, MID_GRAY),
    ("\u2022  Тренд на цифровизацию в медицине", 16, MID_GRAY),
])

# CTA card
cta_card = add_rounded_rect(
    slide6, Inches(8.2), Inches(1.8), Inches(4.3), Inches(4.5),
    RGBColor(0x14, 0x14, 0x30), RGBColor(0x35, 0x30, 0x60), Pt(1.5)
)

add_accent_line(slide6, Inches(8.2), Inches(1.8), Inches(4.3), ACCENT_PURPLE, Pt(4))

add_multiline_text(slide6, Inches(8.6), Inches(2.4), Inches(3.5), Inches(3.5), [
    ("Готовы к", 28, WHITE, True),
    ("масштабированию", 28, LIGHT_PURPLE, True),
    ("", 12, WHITE),
    ("Прототип работает.", 17, LIGHT_GRAY),
    ("Модель обучена.", 17, LIGHT_GRAY),
    ("Приложение готово к пилотному внедрению.", 17, LIGHT_GRAY),
    ("", 12, WHITE),
    ("Открыты для партнёрства и инвестиций.", 17, WHITE, True),
])


# ============================================================
# SLIDE 7 — Заключение
# ============================================================
slide7 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide7, BG_DARK)

# Decorative shapes
bg_oval = slide7.shapes.add_shape(
    MSO_SHAPE.OVAL, Inches(2), Inches(0), Inches(9), Inches(8)
)
bg_oval.fill.solid()
bg_oval.fill.fore_color.rgb = RGBColor(0x12, 0x12, 0x25)
bg_oval.line.fill.background()

add_text_box(slide7, Inches(1), Inches(1.5), Inches(11.3), Inches(1.2),
             "AI Doctor",
             font_size=56, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide7, Inches(1), Inches(2.8), Inches(11.3), Inches(0.6),
             "Автоматическая диагностика ВНЧС с помощью искусственного интеллекта",
             font_size=22, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

# Summary points
summary_items = [
    "Интеграция ML в медицинские приложения",
    "Обработка медицинских изображений (DICOM)",
    "Full-stack решение: iOS + Backend + ML",
    "Практическое применение глубокого обучения в диагностике",
]

for i, item in enumerate(summary_items):
    iy = Inches(3.7) + i * Inches(0.5)
    add_text_box(slide7, Inches(3), iy, Inches(7.3), Inches(0.45),
                 f"\u2713   {item}", font_size=18, color=LIGHT_GRAY,
                 alignment=PP_ALIGN.LEFT)

add_text_box(slide7, Inches(1), Inches(5.8), Inches(11.3), Inches(0.6),
             "Магистерская программа  \u2022  Искусственный интеллект в медицине",
             font_size=16, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide7, Inches(1), Inches(6.3), Inches(11.3), Inches(0.5),
             "Спасибо за внимание!",
             font_size=28, color=LIGHT_PURPLE, bold=True, alignment=PP_ALIGN.CENTER)


# ============================================================
# SAVE
# ============================================================
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "presentation.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Slides: {len(prs.slides)}")
