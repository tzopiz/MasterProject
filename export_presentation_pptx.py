#!/usr/bin/env python3
"""
Экспорт presentation.html в PPTX (каждая секция — один слайд).
Требуется: pip install playwright python-pptx && playwright install chromium
"""
import asyncio
import io
import sys
from pathlib import Path

try:
    from playwright.async_api import async_playwright
    from pptx import Presentation
    from pptx.util import Inches
except ImportError:
    print("Установите зависимости: pip install playwright python-pptx")
    print("Затем: playwright install chromium")
    sys.exit(1)


SLIDE_SELECTORS = [
    (".hero", "Титул"),
    (".ps-section", "Проблема и решение"),
    (".features-section", "Ключевые возможности"),
    (".stats-section", "Метрики"),
    (".wave-section", "Рынок и CTA"),
]


async def main():
    root = Path(__file__).resolve().parent
    html_path = root / "presentation.html"
    pptx_path = root / "presentation.pptx"

    if not html_path.exists():
        print(f"Файл не найден: {html_path}")
        sys.exit(1)

    async with async_playwright() as p:
        browser = await p.chromium.launch()
        page = await browser.new_page(viewport={"width": 1200, "height": 800})
        await page.goto(f"file://{html_path.absolute()}", wait_until="networkidle")

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        for selector, title in SLIDE_SELECTORS:
            try:
                el = await page.query_selector(selector)
                if el:
                    shot = await el.screenshot(type="png")
                    slide_layout = prs.slide_layouts[6]
                    slide = prs.slides.add_slide(slide_layout)
                    pic = slide.shapes.add_picture(
                        io.BytesIO(shot),
                        Inches(0), Inches(0),
                        width=prs.slide_width,
                        height=prs.slide_height,
                    )
            except Exception as e:
                print(f"Пропуск {title}: {e}")

        await browser.close()

    prs.save(str(pptx_path))
    print(f"Готово: {pptx_path}")


if __name__ == "__main__":
    asyncio.run(main())
